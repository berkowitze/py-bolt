# requires: python-pptx, openpyxl (can use pip install for both)
# once you have this, just run `python boltpptx.py`
# if it doesn't work, email eliberkowitz@gmail.com with the error message
# as it stands right now, you must have a spreadsheet with columns for:
## first name              (str)
## last name               (str)
## hometown                (str)
## freshman year unit      (int/str) if not a transfer
## outdoor experience      (int)
## personal adjective      (str)
## medical information     (str)
## dietary restrictions    (str)
## compass estimate        (str) (estimation of N/E/S/W)
## compass input           (str) (what they typed in if can't be classified)
## POC group interest      (int) (1 <= int <= 10)
## Creative group interest (int) (1 <= int <= 10)
## Womens group interest   (int) (1 <= int <= 10)

from pptx import Presentation
import openpyxl
from openpyxl.utils import column_index_from_string as str_to_col

# Set parameters here
# (powerpoint template, spreadsheet location, column letters)
pptx_template  = 'template.pptx'
output_pptx    = 'bolter_cards.pptx'
main_ss        = 'bolters.xlsx'
# Transfers should just have a column with first names and
# a column with last names. Also possible to add a column in the
# main_ss with transfer boolean, but this isn't implemented right now
# (would only be a few lines)
transfers_ss   = 'transfers.xlsx'
fn_col                = 'B'
ln_col                = 'C'
hometown_col          = 'D'
unit_col              = 'E'
people_col            = 'N'
med_col               = 'O'
diet_col              = 'P'
women_interest_col    = 'R'
poc_interest_col      = 'S'
creative_interest_col = 'T'
compass_col           = 'G'
compass_backup        = 'F'
outdoors_col          = 'Q'
adjective_col         = 'K'

# first name column in transfer ss
fn_t_col = 'A'
# last name column in transfer ss
ln_t_col = 'B'

# get transfer first/last names
transfer_xl = openpyxl.load_workbook(transfers_ss).get_active_sheet()
transfer_fns = map(lambda x: x.value, transfer_xl[fn_t_col])
transfer_lns = map(lambda x: x.value, transfer_xl[ln_t_col])

# open the main spreadsheet
name_sheet = openpyxl.load_workbook(main_ss).get_active_sheet()

# open the presentation template and get the layout template
prs = Presentation(open(pptx_template))
layout = prs.slide_layouts[0]
for i, row in enumerate(list(name_sheet.values)[1:]):
    # grab values from the spreadsheet row
    fn        = row[str_to_col(fn_col) - 1]
    ln        = row[str_to_col(ln_col) - 1]
    transfer  = fn in transfer_fns and ln in transfer_lns
    unit      = row[str_to_col(unit_col) - 1]
    women     = row[str_to_col(women_interest_col) - 1] >= 7
    poc       = row[str_to_col(poc_interest_col) - 1] >= 7
    creative  = row[str_to_col(creative_interest_col) - 1] >= 7
    med       = row[str_to_col(med_col) - 1]
    diet      = row[str_to_col(diet_col) - 1]
    people    = row[str_to_col(people_col) - 1]
    compass   = row[str_to_col(compass_col) - 1]
    adjective = row[str_to_col(adjective_col) - 1]
    hometown  = row[str_to_col(hometown_col) - 1]
    outdoors  = row[str_to_col(outdoors_col) - 1]

    group_list = []
    for i, x in enumerate([women, poc, creative]):
        if i == 0 and x:
            group_list.append('Women')
        if i == 1 and x:
            group_list.append('POC')
        if i == 2 and x:
            group_list.append('Creative')

    # construct strings to go on the card
    group_str = 'None' if not group_list else ', '.join(group_list)
    if (isinstance(unit, int) or isinstance(unit, float)):
        unit_str = str(int(unit))
    else:
        unit_str = str(unit)

    if isinstance(compass, unicode) or isinstance(compass, str):
        compass_str = compass
    else:
        compass_str = row[str_to_col(compass_backup) - 1]

    # duplicate slide and get placeholders
    slide = prs.slides.add_slide(layout)
    phs = list(slide.placeholders)

    # populate slide
    phs[0].text = fn + ' ' + ln
    phs[1].text = 'Home: ' + hometown
    phs[2].text = 'Transfer' if transfer else 'Unit ' + unit_str
    phs[3].text = 'People: ' + people
    phs[4].text = 'Diet: ' + diet if diet else ''
    phs[5].text = 'Med: ' + med if med else 'Med: none'
    phs[6].text = 'Experience: ' + str(int(outdoors) if outdoors == int(outdoors) else outdoors)
    phs[7].text = 'Adjective: ' + adjective
    phs[8].text = 'Groups: ' + group_str
    phs[9].text = 'Compass: ' + compass_str

prs.save(output_pptx)

